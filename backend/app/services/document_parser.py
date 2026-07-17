
"""
文档解析与智能切块服务
支持四种切块策略: default(递归), semantic(语义), markdown(MD结构), parent_child(父子)
"""
import re
import os
import hashlib
import csv
import io
import json
import math
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from langchain_text_splitters import RecursiveCharacterTextSplitter
from app.core.config import get_settings

settings = get_settings()

SUPPORTED_EXTS = {".pdf", ".pptx", ".ppt", ".md", ".markdown", ".txt", ".docx", ".doc", ".csv"}


# ============================================================
# 文本提取
# ============================================================

def extract_text_from_pdf(file_path: str) -> str:
    try:
        import pdfplumber
        # 提前检测空文件
        if os.path.getsize(file_path) == 0:
            raise ValueError("文档内容为空（文件大小为 0 字节），请检查文件是否损坏")
        parts = []
        with pdfplumber.open(file_path, password="") as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)
        text = "\n".join(parts)
        if not text.strip():
            raise ValueError("PDF 无法提取文字，可能为扫描件（纯图片）或加密文档，请检查文件")
        return text
    except ValueError:
        raise
    except Exception as e:
        msg = str(e)
        if "password" in msg.lower() or "encrypt" in msg.lower():
            raise ValueError("PDF 文档已加密，请先解密后再上传")
        raise ValueError(f"PDF 解析失败: {msg}")


def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception:
        raise ValueError("PPT/PPTX 解析失败")


def extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        return "\n".join(parts)
    except Exception:
        raise ValueError("Word 文档解析失败")


def extract_text_from_md(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_text_from_csv(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        rows = list(reader)
    if not rows:
        return ""
    header = rows[0]
    col_count = len(header)
    parts = [f"### CSV表格 (共 {col_count} 列, {len(rows)-1} 行数据)"]
    parts.append(" | ".join(header))
    parts.append(" | ".join(["---"] * col_count))
    for row in rows[1:]:
        padded = row + [""] * (col_count - len(row))
        parts.append(" | ".join(padded[:col_count]))
    return "\n".join(parts)


def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_path)
    elif ext in (".md", ".markdown"):
        return extract_text_from_md(file_path)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    else:
        return extract_text_from_txt(file_path)


# ============================================================
# 通用工具函数
# ============================================================

def _make_chunk(content: str, idx: int, chunk_type: str = "child",
                parent_id: Optional[int] = None,
                metadata: Optional[dict] = None) -> Dict[str, Any]:
    """构建标准 chunk 字典"""
    m = metadata or {}
    m["chunk_type"] = chunk_type
    return {
        "content": content,
        "chunk_index": idx,
        "chunk_type": chunk_type,
        "parent_chunk_id": parent_id,
        "chunk_hash": hashlib.sha256(content.encode("utf-8")).hexdigest(),
        "metadata": json.dumps(m, ensure_ascii=False),
    }


def _split_sentences(text: str) -> List[str]:
    """中英文通用分句"""
    # 统一换行处理
    text = re.sub(r'\n+', ' ', text)
    # 按句末标点切分：。！？! ? . 但排除小数点
    pattern = r'(?<=[。！？!?])(?=[^"\'」』\)）])|(?<=[.](?<![0-9]\.[^.]))(?=\s+[A-ZÀ-ÖØ-Þ])'
    raw = re.split(pattern, text)
    sentences = []
    for s in raw:
        s = s.strip()
        if s and len(s) > 1:  # 过滤空句和单字符
            sentences.append(s)
    return sentences


def _estimate_chinese_chars(text: str) -> int:
    """估算中文字符数（用于 chunk_size 换算）"""
    cn = sum(1 for c in text if '\u4e00' <= c <= '\u9fff' or '\u3400' <= c <= '\u4dbf')
    en = len(text) - cn
    # 1 中文字符 ≈ 1.5 个 token，1 英文词 ≈ 1.3 个 token
    return int(cn + en * 0.4)


# ============================================================
# 策略 1: Default — 增强递归切块
# ============================================================

def split_text(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
    """
    递归字符切块 (默认策略)
    针对中英文混合文档优化：
    - 优先按段落分隔符切分
    - 加入中文句号、问号、感叹号作为分隔点
    """
    cs = chunk_size or settings.chunk_size
    co = chunk_overlap or settings.chunk_overlap

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=cs,
        chunk_overlap=co,
        length_function=len,
        separators=[
            "\n\n\n",         # 多空行（节分隔）
            "\n\n",           # 段落
            "\n",             # 行
            "。",             # 中文句号
            "！", "!",        # 感叹号
            "？", "?",        # 问号
            "；", ";",        # 分号
            "，", ",",        # 逗号
            ". ",             # 英文句号+空格
            " ",              # 空格
            "",               # 字符级
        ],
        keep_separator=True,
    )

    raw_chunks = splitter.split_text(text)
    chunks = []
    for i, c in enumerate(raw_chunks):
        chunks.append(_make_chunk(
            content=c.strip(),
            idx=i,
            chunk_type="child",
            metadata={"method": "default", "char_count": len(c)}
        ))
    return chunks


# ============================================================
# 策略 2: Semantic — 真正的语义切块
# ============================================================

def _compute_embedding(texts: List[str]) -> List[List[float]]:
    """计算文本的 embedding 向量"""
    try:
        import asyncio
        import concurrent.futures
        from app.services.llm_service import get_embeddings
        # asyncio.run() 不能在已运行的事件循环中调用（如 FastAPI）
        # 所以在独立线程中执行
        def _run():
            return asyncio.run(get_embeddings(texts))
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            return pool.submit(_run).result(timeout=60)
    except Exception:
        return None


def split_text_semantic(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
    """
    语义切块策略

    核心原理:
    1. 将文本拆分为句子
    2. 逐句计算 embedding
    3. 计算相邻句子组的语义相似度
    4. 当相似度骤降时（话题转换），在此处切分
    5. 将语义相近的句子合并为一个 chunk

    这种方法能自动识别文档中的主题边界，产生比固定长度更自然的切分。
    """
    cs = chunk_size or settings.chunk_size

    # 第一步：分句
    sentences = _split_sentences(text)
    if len(sentences) <= 1:
        return split_text(text, cs)  # 降级为默认策略

    # 第二步：计算 embedding
    embeddings = _compute_embedding(sentences)
    if embeddings is None or len(embeddings) != len(sentences):
        return split_text(text, cs)  # embedding 失败时降级

    # 第三步：计算相邻句子相似度
    def cosine(a: List[float], b: List[float]) -> float:
        dot = sum(x * y for x, y in zip(a, b))
        norm_a = math.sqrt(sum(x * x for x in a))
        norm_b = math.sqrt(sum(x * x for x in b))
        if norm_a == 0 or norm_b == 0:
            return 0.0
        return dot / (norm_a * norm_b)

    # 第四步：自适应阈值检测断点
    # 计算所有相邻句子的相似度
    similarities = []
    for i in range(len(sentences) - 1):
        sim = cosine(embeddings[i], embeddings[i + 1])
        similarities.append(sim)

    if not similarities:
        return split_text(text, cs)

    # 使用百分位阈值：低于 25% 分位数的相似度视为断点
    sorted_sims = sorted(similarities)
    p25 = sorted_sims[int(len(sorted_sims) * 0.25)]
    # 阈值不低于 0.5（如果整体都很相似，不强制切分）
    threshold = max(p25, 0.5)

    # 标记断点位置
    breakpoints = []
    for i, sim in enumerate(similarities):
        if sim < threshold:
            breakpoints.append(i + 1)  # 在第 i+1 句后切断

    # 第五步：按断点合并句子为 chunks
    chunks = []
    current_group = [sentences[0]]
    current_len = len(sentences[0])

    for i in range(1, len(sentences)):
        sen_len = len(sentences[i])

        # 如果累积长度超过 chunk_size 或遇到断点，开始新 chunk
        if current_len + sen_len > cs and current_group:
            chunks.append("".join(current_group))
            current_group = []
            current_len = 0

        if i in breakpoints and current_group:
            chunks.append("".join(current_group))
            current_group = []
            current_len = 0

        current_group.append(sentences[i])
        current_len += sen_len

    if current_group:
        chunks.append("".join(current_group))

    # 第六步：合并过小的 chunk
    merged_chunks = []
    buffer = ""
    min_chunk_size = cs // 3  # 最小 chunk 大小为目标的 1/3
    for c in chunks:
        if buffer and len(buffer) + len(c) <= cs * 1.2:
            buffer += "\n\n" + c
        elif buffer:
            merged_chunks.append(buffer.strip())
            buffer = c
        else:
            buffer = c
    if buffer:
        merged_chunks.append(buffer.strip())

    # 构建结果
    result = []
    for i, c in enumerate(merged_chunks):
        result.append(_make_chunk(
            content=c,
            idx=i,
            chunk_type="child",
            metadata={"method": "semantic", "char_count": len(c),
                       "breakpoint_count": len(breakpoints)}
        ))
    return result


# ============================================================
# 策略 3: Markdown — 结构化切块
# ============================================================

def _parse_md_headings(lines: List[str]) -> List[Dict[str, Any]]:
    """
    解析 Markdown 文本的结构：
    返回 [{type, level, title, line_start, line_end, children_start, children_end, code_block}]
    """
    sections = []
    in_code_block = False
    heading_pattern = re.compile(r'^(#{1,6})\s+(.+)$')
    table_sep = re.compile(r'^\|[\s\-:|]+\|?$')

    i = 0
    while i < len(lines):
        line = lines[i]

        # 检测代码块边界
        code_fence = line.strip().startswith("```") or line.strip().startswith("~~~")
        if code_fence:
            in_code_block = not in_code_block
            i += 1
            continue

        # 跳过表格分隔行
        if table_sep.match(line.strip()):
            i += 1
            continue

        # 检测标题
        m = heading_pattern.match(line)
        if m and not in_code_block:
            level = len(m.group(1))
            title = m.group(2).strip()
            sections.append({
                "type": "heading",
                "level": level,
                "title": title,
                "line": i,
            })
        i += 1

    # 添加每个 section 的结束行
    for j, sec in enumerate(sections):
        if j + 1 < len(sections):
            sec["line_end"] = sections[j + 1]["line"] - 1
        else:
            sec["line_end"] = len(lines) - 1

    return sections


def _build_heading_path(sections: List[Dict], idx: int) -> str:
    """构建从文档根到当前标题的层级路径（只包含严格上级标题，排除同级兄弟）"""
    path = [sections[idx]["title"]]
    current_level = sections[idx]["level"]
    # 从当前位置向前回溯，只收集严格更高级别的标题作为父级
    for j in range(idx - 1, -1, -1):
        sec = sections[j]
        if sec["level"] < current_level:
            path.insert(0, sec["title"])
            current_level = sec["level"]
    return " > ".join(path)


def split_text_markdown(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
    """
    Markdown 结构感知切块策略

    核心原理:
    1. 解析所有标题层级（H1-H6）
    2. 按标题将文档分割为逻辑 section
    3. 每个 chunk 保留其标题路径作为上下文
    4. 过大的 section 内部再按段落切分，并带上父标题
    5. 保护代码块不被切断
    """
    cs = chunk_size or settings.chunk_size
    lines = text.split("\n")
    sections = _parse_md_headings(lines)

    # 如果没有标题结构，降级为默认策略
    if not sections:
        return split_text(text, cs)

    # 提取每个 section 的正文
    in_code_block = False
    chunks = []
    chunk_idx = 0

    for sec_idx, sec in enumerate(sections):
        heading_path = _build_heading_path(sections, sec_idx)
        start = sec["line"] + 1  # 正文从标题下一行开始
        end = sec["line_end"] + 1  # 包含结束行

        # 收集 section 正文
        body_lines = []
        for li in range(start, min(end, len(lines))):
            line = lines[li]
            if line.strip().startswith("```") or line.strip().startswith("~~~"):
                in_code_block = not in_code_block
            body_lines.append(line)

        body_text = "\n".join(body_lines).strip()
        if not body_text:
            continue

        heading_prefix = f"# {heading_path}"

        # 如果正文长度在 chunk_size 内，直接作为一个 chunk
        if len(body_text) <= cs:
            full_text = f"{heading_prefix}\n\n{body_text}"
            chunks.append(_make_chunk(
                content=full_text,
                idx=chunk_idx,
                chunk_type="child",
                metadata={
                    "method": "markdown",
                    "heading": heading_path,
                    "heading_level": sec["level"],
                    "char_count": len(full_text),
                }
            ))
            chunk_idx += 1
        else:
            # 正文太长，用递归切分再拆，每个子 chunk 带标题前缀
            sub_texts = body_text.split("\n\n")
            buffer = ""
            for sub in sub_texts:
                sub = sub.strip()
                if not sub:
                    continue
                candidate = f"{heading_prefix}\n\n{buffer}\n\n{sub}" if buffer else f"{heading_prefix}\n\n{sub}"
                if len(candidate) > cs and buffer:
                    chunks.append(_make_chunk(
                        content=f"{heading_prefix}\n\n{buffer.strip()}",
                        idx=chunk_idx,
                        chunk_type="child",
                        metadata={
                            "method": "markdown",
                            "heading": heading_path,
                            "heading_level": sec["level"],
                            "char_count": len(buffer.strip()),
                        }
                    ))
                    chunk_idx += 1
                    buffer = sub
                else:
                    buffer = candidate.replace(f"{heading_prefix}\n\n", "", 1) if buffer else sub

            if buffer:
                chunks.append(_make_chunk(
                    content=f"{heading_prefix}\n\n{buffer.strip()}",
                    idx=chunk_idx,
                    chunk_type="child",
                    metadata={
                        "method": "markdown",
                        "heading": heading_path,
                        "heading_level": sec["level"],
                        "char_count": len(buffer.strip()),
                    }
                ))
                chunk_idx += 1

    # 合并过小的相邻 chunk（同层级）
    merged = []
    buffer = None
    for ch in chunks:
        if buffer is None:
            buffer = ch
        elif (ch.get("metadata") and buffer.get("metadata")
              and json.loads(ch["metadata"]).get("heading") == json.loads(buffer["metadata"]).get("heading")
              and len(buffer["content"]) + len(ch["content"]) <= cs * 1.2):
            buffer["content"] += "\n\n" + ch["content"].replace(f"# {json.loads(ch['metadata']).get('heading', '')}\n\n", "", 1)
            buffer["chunk_hash"] = hashlib.sha256(buffer["content"].encode("utf-8")).hexdigest()
            m = json.loads(buffer["metadata"])
            m["char_count"] = len(buffer["content"])
            buffer["metadata"] = json.dumps(m, ensure_ascii=False)
        else:
            merged.append(buffer)
            buffer = ch
    if buffer:
        merged.append(buffer)

    for i, ch in enumerate(merged):
        ch["chunk_index"] = i

    return merged


# ============================================================
# 策略 4: Parent-Child — 父子切块 + 两阶段检索
# ============================================================

def split_text_parent_child(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
    """
    父子切块策略

    核心原理:
    1. Parent chunks（父块）：较大粒度，提供完整上下文（如整个 section）
    2. Child chunks（子块）：较小粒度，用于精确匹配（如单个段落）
    3. 两阶段检索：用问题匹配子块 → 取匹配子块的父块作为最终上下文

    这种策略兼顾了检索精度和上下文完整性。
    """
    cs = chunk_size or settings.chunk_size
    parent_size = cs * 3  # 父块是子块的 3 倍
    child_size = cs

    # 先尝试按 Markdown 结构解析（如果有标题）
    lines = text.split("\n")
    sections = _parse_md_headings(lines)

    all_chunks = []
    parent_idx = 0
    child_base_id = 100000  # 虚拟 ID，由数据库生成真实 ID

    def _add_parent_child(parent_text: str, child_texts: List[str], heading_path: str = "", level: int = 0):
        nonlocal parent_idx, child_base_id
        parent_id = parent_idx + 1
        meta = {"method": "parent_child", "char_count": len(parent_text)}
        if heading_path:
            meta["heading"] = heading_path
            meta["heading_level"] = level

        # 添加父块
        all_chunks.append(_make_chunk(
            content=parent_text, idx=parent_idx,
            chunk_type="parent", parent_id=None, metadata=meta.copy()
        ))
        parent_idx += 1

        # 添加子块
        for ci, ct in enumerate(child_texts):
            child_meta = {**meta, "char_count": len(ct)}
            all_chunks.append(_make_chunk(
                content=ct, idx=child_base_id + ci,
                chunk_type="child", parent_id=parent_id, metadata=child_meta
            ))
        child_base_id += len(child_texts)

    if sections:
        for sec in sections:
            heading_path = _build_heading_path(sections, sections.index(sec))
            start = sec["line"] + 1
            end = sec["line_end"] + 1
            body = "\n".join(lines[start:min(end, len(lines))]).strip()
            if not body:
                continue

            heading_prefix = f"# {heading_path}\n\n"
            parent_text = heading_prefix + body

            # 子块：按段落拆分
            paragraphs = [p.strip() for p in body.split("\n\n") if p.strip()]
            child_texts = [heading_prefix + p for p in paragraphs]

            # 如果父块太大，进一步拆分子父块
            if len(parent_text) > parent_size * 1.5:
                # 对父块再递归拆分
                sub_parents = []
                sub_children = []
                buf = []
                buf_len = 0
                for p in paragraphs:
                    if buf_len + len(p) > parent_size and buf:
                        sub_parents.append("\n\n".join(buf))
                        sub_children.append(buf)
                        buf = []
                        buf_len = 0
                    buf.append(p)
                    buf_len += len(p)
                if buf:
                    sub_parents.append("\n\n".join(buf))
                    sub_children.append(buf)

                for sp, sc in zip(sub_parents, sub_children):
                    _add_parent_child(
                        f"{heading_prefix}{sp}",
                        [f"{heading_prefix}{c}" for c in sc],
                        heading_path, sec["level"]
                    )
            else:
                _add_parent_child(parent_text, child_texts, heading_path, sec["level"])
    else:
        # 无标题结构：用默认切分生成父子
        splitter_large = RecursiveCharacterTextSplitter(
            chunk_size=parent_size, chunk_overlap=100,
            separators=["\n\n", "\n", "。", "！", "？", ". ", " ", ""]
        )
        splitter_small = RecursiveCharacterTextSplitter(
            chunk_size=child_size, chunk_overlap=50,
            separators=["\n\n", "\n", "。", "！", "？", ". ", " ", ""]
        )
        parent_chunks = splitter_large.split_text(text)
        for pc in parent_chunks:
            child_sub = splitter_small.split_text(pc)
            _add_parent_child(pc.strip(), [c.strip() for c in child_sub if c.strip()])

    return all_chunks



# ============================================================
# CSV 智能切块 — 保护表格结构不被切碎
# ============================================================

def split_csv_table(text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[Dict[str, Any]]:
    """
    CSV 表格智能切块

    核心原则：不破坏表格结构，切分粒度可控
    - 表格 < chunk_size → 单 chunk（完整表格）
    - 表格 > chunk_size → 按行分组切块，每个 chunk 重复表头
    """
    cs = chunk_size or settings.chunk_size
    lines = text.split("\n")
    chunks = []
    chunk_idx = 0

    # 解析表格结构：找表头行和分隔行
    header_line = -1
    sep_line = -1
    for i, line in enumerate(lines):
        stripped = line.strip()
        if not stripped:
            continue
        if "---" in stripped and "|" in stripped:
            sep_line = i
            header_line = i - 1
            break
        if "|" in stripped and header_line < 0:
            header_line = i

    # 不是标准 Markdown 表格，降级为普通切分
    if header_line < 0:
        return split_text(text, cs, chunk_overlap)

    header = lines[header_line].strip()
    table_start = header_line
    data_start = sep_line + 1 if sep_line > header_line else header_line + 1
    data_lines = [l for l in lines[data_start:] if l.strip()]

    # 小表格：一个 chunk 搞定
    full_len = len(header) + sum(len(l) for l in data_lines)
    if full_len <= cs and len(data_lines) <= 50:
        chunks.append({
            "chunk_index": chunk_idx,
            "content": text.strip(),
            "chunk_hash": hashlib.md5(text.strip().encode()).hexdigest(),
            "metadata": json.dumps({"type": "csv_table", "columns": _csv_column_names(header), "rows": len(data_lines), "split": "full"}, ensure_ascii=False),
        })
        return chunks

    # 大表格：按行分组切块，每个块带表头
    rows_per_chunk = max(1, (cs - len(header) - 50) // max(1, min(len(l) for l in data_lines) if data_lines else 50))
    header_block = lines[table_start:data_start]  # 表头+分隔行

    for i in range(0, len(data_lines), rows_per_chunk):
        batch = header_block + data_lines[i:i + rows_per_chunk]
        chunk_text = "\n".join(batch).strip()
        chunks.append({
            "chunk_index": chunk_idx,
            "content": chunk_text,
            "chunk_hash": hashlib.md5(chunk_text.encode()).hexdigest(),
            "metadata": json.dumps({
                "type": "csv_table",
                "columns": _csv_column_names(header),
                "row_start": i + 1,
                "row_end": min(i + rows_per_chunk, len(data_lines)),
                "total_rows": len(data_lines),
            }, ensure_ascii=False),
        })
        chunk_idx += 1

    return chunks


def _csv_column_names(header_line: str) -> List[str]:
    """提取 CSV 表格的列名"""
    return [c.strip() for c in header_line.split("|") if c.strip() and c.strip() != ""]

# ============================================================
# 策略映射 & 文档处理入口
# ============================================================

CHUNK_METHODS = {
    "default": split_text,
    "semantic": split_text_semantic,
    "markdown": split_text_markdown,
    "parent_child": split_text_parent_child,
    "csv_table": split_csv_table,
}


def process_document(file_path: str, chunk_size: int = None,
                     chunk_overlap: int = None,
                     chunk_method: str = "default") -> dict:
    """
    统一的文档处理入口

    返回:
    {
        "filename": str,
        "full_text": str,
        "chunks": List[chunk_dict],
        "chunk_count": int,
        "method": str,
        "parent_count": int,
        "child_count": int,
    }
    """
    ext = Path(file_path).suffix.lower()
    text = extract_text(file_path)
    if not text or not text.strip():
        raise ValueError("文档内容为空，无法提取文字。请确认文件包含可读文本")

    # CSV 自动使用表格感知切块策略，保护表格结构
    if ext == ".csv":
        chunk_method = "csv_table"

    method_fn = CHUNK_METHODS.get(chunk_method, split_text)
    chunks = method_fn(text, chunk_size, chunk_overlap)

    parent_count = sum(1 for c in chunks if c.get("chunk_type") == "parent")
    child_count = sum(1 for c in chunks if c.get("chunk_type") == "child")

    return {
        "filename": Path(file_path).name,
        "full_text": text,
        "chunks": chunks,
        "chunk_count": len(chunks),
        "method": chunk_method,
        "parent_count": parent_count,
        "child_count": child_count,
    }
